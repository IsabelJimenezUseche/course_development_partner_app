import hashlib
import json
import os
import re
import socket
import subprocess
import sys
import zipfile
from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from fastapi import HTTPException
from fastapi.testclient import TestClient
from pptx import Presentation
import server

from server import (
    _apply_auto_decision,
    _artifact_spec_markdown,
    _collaboration_mode_instruction,
    _extract_decision,
    _extract_artifact_spec,
    _extract_state_file,
    _enforce_professor_question_limit,
    _recover_professor_clarification,
    _hide_incomplete_internal_payloads,
    _infer_skill_profile,
    _hashed_text_vector,
    _load_skill_runtime,
    _persist_exchange,
    _prepare_professor_display_content,
    _render_markdown,
    _sanitize_professor_decision,
    _unwrap_markdown_documents,
    _validate_project_state,
    _vector_cosine,
    _write_state_file,
    app,
    find_available_port,
    get_settings,
    ChatMessage,
    PROFESSOR_INTERACTION_CONTRACT,
    STATE_FILE_VALIDATORS,
)


def _skill_fixture(name: str, filename: str) -> str:
    """Read a fixture shipped with the installed skill, so tests track the real thing."""
    return (
        get_settings().skill_dir.parent / "tests" / "fixtures" / name / filename
    ).read_text(encoding="utf-8")


client = TestClient(app)


def test_local_vector_retrieval_favors_relevant_document_chunk():
    query = _hashed_text_vector("scaled dot product attention queries keys values")
    relevant = _hashed_text_vector(
        "Scaled dot product attention compares queries with keys and combines the values."
    )
    distractor = _hashed_text_vector(
        "The course registration calendar lists deadlines for adding and dropping classes."
    )

    assert _vector_cosine(query, relevant) > _vector_cosine(query, distractor)


def test_health_does_not_expose_api_key(monkeypatch):
    secret = "test-secret-that-must-not-appear"
    monkeypatch.setenv("PURDUE_GENAI_API_KEY", secret)

    response = client.get("/health")

    assert response.status_code == 200
    assert response.json()["api_key_configured"] is True
    assert secret not in response.text


def test_safe_config_does_not_expose_api_key(monkeypatch):
    secret = "another-test-secret"
    monkeypatch.setenv("PURDUE_GENAI_API_KEY", secret)

    response = client.get("/api/config")

    assert response.status_code == 200
    assert response.json()["api_key_configured"] is True
    assert secret not in response.text


def test_root_serves_faculty_ui():
    response = client.get("/")

    assert response.status_code == 200
    assert "Course Development Partner" in response.text
    assert "Design decisions that stay connected" in response.text
    assert "Generated artifacts" in response.text


def test_skill_runtime_loads_actual_skill_files():
    response = client.get("/api/skill", params={"profile": "assessment"})

    assert response.status_code == 200
    runtime = response.json()
    assert runtime["name"] == "course-development-partner"
    assert runtime["loaded_files"] == [
        "SKILL.md",
        "references/assessment-quality.md",
        "references/artifact-patterns.md",
    ]
    assert len(runtime["fingerprint"]) == 16


def test_stem_profile_loads_renamed_authenticity_reference():
    _, runtime = _load_skill_runtime(get_settings(), "stem")

    assert runtime["loaded_files"] == [
        "SKILL.md",
        "references/stem-authenticity.md",
        "references/artifact-patterns.md",
    ]


def test_retired_engineering_profile_still_resolves_to_stem():
    _, runtime = _load_skill_runtime(get_settings(), "engineering")

    assert runtime["profile"] == "stem"
    assert "references/stem-authenticity.md" in runtime["loaded_files"]


def test_stem_profile_routes_beyond_engineering_wording():
    for text in ("Plan a laboratory session", "Draft a computing lab on uncertainty"):
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == "stem"


def test_framework_profile_loads_crosswalk_and_recording_state():
    """The skill update added framework-crosswalk.md; a named framework must reach it."""
    _, runtime = _load_skill_runtime(get_settings(), "framework")

    assert runtime["loaded_files"] == [
        "SKILL.md",
        "references/framework-crosswalk.md",
        "references/evidence-informed-design.md",
    ]
    assert runtime["loaded_assets"] == [
        "assets/design-log.md",
        "assets/source-register.md",
    ]


def test_framework_profile_is_available_through_the_api():
    response = client.get("/api/skill", params={"profile": "framework"})

    assert response.status_code == 200
    assert "references/framework-crosswalk.md" in response.json()["loaded_files"]
    assert client.get("/api/skill", params={"profile": "not-a-profile"}).status_code == 400


def test_named_framework_routes_to_the_crosswalk_profile():
    for text in (
        "Map my course outcomes to the ABET student outcomes",
        "How do these fit Bloom's taxonomy levels?",
        "We teach with POGIL and want to keep the roles working",
        "Our program is preparing for accreditation next year",
        "I want to restructure this as a flipped classroom",
    ):
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == "framework"


def test_framework_words_inside_other_words_do_not_reroute():
    """'diabetes' contains 'abet' and 'algal blooms' nearly contains "bloom's"."""
    for text in (
        "Plan a unit on diabetes for nursing students",
        "Cover seasonal algal blooms in the ecology unit",
    ):
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == "establish"


def test_framework_mention_defers_to_the_specific_task_route():
    """A rubric that cites ABET is rubric work; the crosswalk is for pure mapping asks."""
    profile = _infer_skill_profile(
        [ChatMessage(role="user", content="Draft a rubric for ABET outcome 2")]
    )

    assert profile == "assessment"


def test_named_grading_schemes_route_to_assessment_not_framework():
    """Every named scheme contains "grading", so the assessment route claims it first —
    which is where the crosswalk's own §5 sends grading schemes. Pin that deliberately
    so nobody re-adds the scheme names to FRAMEWORK_SIGNAL as reachable patterns."""
    for text in (
        "Help me set up specifications grading for this course",
        "I want to move to ungrading next semester",
        "Explain standards-based grading to my TAs",
    ):
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == "assessment"


def test_syllabus_and_delivery_adaptation_route_to_course_profile():
    """The skill update added syllabi and online/hybrid adaptation to its scope."""
    for text in (
        "Draft a syllabus for my spring section",
        "Move my seminar to hybrid delivery next term",
        "Plan the online adaptation of my lecture course",
    ):
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == "course"


def test_ada_keyword_requires_a_whole_word():
    profile = _infer_skill_profile(
        [ChatMessage(role="user", content="Is this handout ADA compliant?")]
    )
    assert profile == "accessibility"

    # "Adapt" must not be read as the ADA statute.
    profile = _infer_skill_profile(
        [ChatMessage(role="user", content="Adapt my seminar for hybrid delivery")]
    )
    assert profile == "course"


def test_design_profile_supplies_the_verified_bibliography():
    """SKILL.md forbids citing from memory, so the verified entries must be in the prompt."""
    prompt, runtime = _load_skill_runtime(get_settings(), "design")

    assert "references/bibliography.md" in runtime["loaded_files"]
    # The installed file's own content proves the entries reach the prompt verbatim,
    # without pinning this test to any particular bibliography entry.
    bibliography = (
        get_settings().skill_dir / "references" / "bibliography.md"
    ).read_text(encoding="utf-8")
    assert bibliography.strip() in prompt


def test_artifact_profiles_load_the_visual_design_mode_rule():
    """The skill update turned palette selection into a mode rule (visual-design.md §3)
    that artifact-patterns.md defers to; producing visuals without the file in the
    prompt invites improvised colors, which the rule forbids in every mode."""
    for profile in ("artifact", "accessibility"):
        _, runtime = _load_skill_runtime(get_settings(), profile)
        assert "references/visual-design.md" in runtime["loaded_files"], profile


def test_artifact_tool_colors_come_from_the_skill_example_palette():
    """visual-design.md rules out improvised colors: every value the deterministic
    renderer uses must be a role from the example palette table, whose verified
    contrast pairs are the recorded evidence."""
    import artifact_tools

    reference = (
        get_settings().skill_dir / "references" / "visual-design.md"
    ).read_text(encoding="utf-8")
    palette = set(re.findall(r"`#([0-9A-Fa-f]{6})`", reference))
    used = {
        value
        for name, value in vars(artifact_tools).items()
        if name.isupper() and isinstance(value, str) and re.fullmatch(r"[0-9A-Fa-f]{6}", value)
    }
    assert used <= palette, f"colors outside the example palette: {sorted(used - palette)}"


def test_every_profile_route_names_installed_skill_files():
    """A skill rename or removal must fail here, not as a 503 in front of a professor."""
    settings = get_settings()
    for profile, references in server.SKILL_REFERENCE_ROUTES.items():
        for reference in references:
            assert (settings.skill_dir / "references" / reference).is_file(), (
                f"profile {profile} routes references/{reference}, which is not installed"
            )
        for asset in server.SKILL_ASSET_ROUTES.get(profile, []):
            assert (settings.skill_dir / "assets" / asset).is_file(), (
                f"profile {profile} routes assets/{asset}, which is not installed"
            )


def test_skill_references_are_routed_or_deliberately_excluded():
    """When a skill update adds a reference (bibliography.md and framework-crosswalk.md
    arrived this way), the app must either route it into a profile or record here why
    it stays out of the prompt. An unlisted reference is an undecided one."""
    skill_md = (get_settings().skill_dir / "SKILL.md").read_text(encoding="utf-8")
    linked = set(re.findall(r"references/([a-z0-9-]+\.md)", skill_md))
    routed = {
        reference
        for references in server.SKILL_REFERENCE_ROUTES.values()
        for reference in references
    }
    deliberately_excluded = {
        # The app pins its own five tools; connector and tool selection never reach the model.
        "tool-routing.md",
        "mcp-capability-contracts.md",
        # Office production is delegated to the local deterministic artifact tool,
        # which renders with the example palette from visual-design.md; that reference
        # itself is routed so previews name the palette per its §3 mode rule.
        "rich-artifact-production.md",
        # The asset templates are supplied verbatim instead; they carry the schema.
        "state-contract.md",
        # External-search protocol; the app has no web retrieval, and the verified
        # bibliography rides with the design profile instead.
        "evidence-source-protocol.md",
        # End-to-end narrative, too long for a per-turn prompt.
        "worked-example.md",
    }
    assert not deliberately_excluded & routed, "an excluded reference became routed; prune it"
    unrouted = linked - routed - deliberately_excluded
    assert not unrouted, f"skill references neither routed nor excluded: {sorted(unrouted)}"


def test_state_file_registry_matches_the_installed_assets_exactly():
    """The registry is the write allow-list: a state file the skill ships but the app
    does not register can never reach disk, and that silent refusal is exactly what a
    skill update would produce. Fail loudly here instead."""
    installed = {path.name for path in (get_settings().skill_dir / "assets").glob("*.md")}

    assert installed == set(STATE_FILE_VALIDATORS)


def test_registered_validators_are_installed():
    scripts_dir = get_settings().skill_dir / "scripts"
    expected = {name for name in STATE_FILE_VALIDATORS.values() if name}
    expected.add(server.PROJECT_VALIDATOR)

    for script in sorted(expected):
        assert (scripts_dir / script).is_file(), f"scripts/{script} is not installed"


def test_work_dir_defaults_to_the_app_directory(monkeypatch):
    monkeypatch.delenv("APP_WORK_DIR", raising=False)
    monkeypatch.delenv("APP_DATA_DIR", raising=False)

    settings = get_settings()

    assert settings.work_dir == server.APP_DIR
    assert settings.data_dir == server.APP_DIR / "data"


def test_relative_data_dir_resolves_against_the_work_dir(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_WORK_DIR", str(tmp_path))
    monkeypatch.setenv("APP_DATA_DIR", "./elsewhere")

    settings = get_settings()

    assert settings.work_dir == tmp_path.resolve()
    assert settings.data_dir == (tmp_path / "elsewhere").resolve()


def test_absolute_data_dir_ignores_the_work_dir(monkeypatch, tmp_path):
    absolute = tmp_path / "absolute-data"
    monkeypatch.setenv("APP_WORK_DIR", str(tmp_path / "somewhere-else"))
    monkeypatch.setenv("APP_DATA_DIR", str(absolute))

    assert get_settings().data_dir == absolute.resolve()


def test_work_dir_expands_a_home_relative_path(monkeypatch):
    monkeypatch.setenv("APP_WORK_DIR", "~/course-partner-workspace")

    assert get_settings().work_dir == (
        Path.home() / "course-partner-workspace"
    ).resolve()


def test_data_dir_outside_the_app_install_is_created_and_usable(monkeypatch, tmp_path):
    """A data directory on another path must work end to end, not just resolve."""
    monkeypatch.setenv("APP_WORK_DIR", str(tmp_path))
    monkeypatch.setenv("APP_DATA_DIR", "./relocated")

    response = client.post("/api/projects", json={"name": "Relocated data"})

    assert response.status_code == 201
    assert (tmp_path / "relocated" / "workspace.sqlite3").is_file()


def test_skill_runtime_supplies_asset_templates_for_the_profile():
    prompt, runtime = _load_skill_runtime(get_settings(), "assessment")

    assert runtime["loaded_assets"] == [
        "assets/assessment-blueprint.md",
        "assets/alignment-map.md",
    ]
    # The template body must actually reach the model, not just be named.
    assert "Portable state templates" in prompt
    assert "Assessment Blueprint" in prompt


def test_skill_runtime_keeps_internal_production_language_away_from_professor():
    prompt, _ = _load_skill_runtime(get_settings(), "assessment")

    assert PROFESSOR_INTERACTION_CONTRACT in prompt
    assert "Ask no more than three short, direct questions" in prompt
    assert "Never ask the educator for an outcome ID" in prompt
    assert "whether it is practice or graded" in prompt


def test_state_file_contract_lists_only_known_files():
    prompt, _ = _load_skill_runtime(get_settings(), "establish")

    assert "state_file" in prompt
    for filename in STATE_FILE_VALIDATORS:
        assert filename in prompt


def test_extract_state_file_rejects_unknown_file_names():
    content, spec = _extract_state_file(
        'Here you go.\n\n```state_file\n{"file":"../../etc/passwd","content":"x"}\n```'
    )

    assert spec is None
    assert "state_file" in content  # left intact rather than silently written


def test_extract_state_file_recovers_when_gpt_oss_drops_the_fence():
    """The target model reliably emits the payload but often omits the fence."""
    bare = '{"file":"design-log.md","content":"# Design Log\\n\\n- Schema version: 1.0\\n"}'

    content, spec = _extract_state_file(bare)

    assert spec["file"] == "design-log.md"
    assert "Schema version" in spec["content"]
    # A reply that was nothing but the payload still needs display text.
    assert content and "design-log.md" in content


def test_extract_state_file_ignores_json_that_is_not_a_state_file():
    for text in (
        "Plain prose with no JSON.",
        'Example: ```json\n{"file":"whatever.txt","content":"x"}\n```',
        '```json\n{"kind":"slides","title":"Not a state file"}\n```',
    ):
        _, spec = _extract_state_file(text)
        assert spec is None


def test_extract_state_file_returns_known_file():
    content, spec = _extract_state_file(
        'Saved.\n\n```state_file\n{"file":"design-log.md","content":"# Design Log\\n"}\n```'
    )

    assert spec["file"] == "design-log.md"
    assert content == "Saved."


def test_extract_state_file_recovers_known_markdown_template_fence():
    content, spec = _extract_state_file(
        "The lesson state is ready.\n\n```markdown\n"
        "# Lesson Storyboard\n\n- Schema version: 1.0\n\n## Sequence\n\n"
        "| Time | Purpose |\n|---|---|\n| 0-5 | Predict |\n```"
    )

    assert content == "The lesson state is ready."
    assert spec["file"] == "lesson-storyboard.md"
    assert spec["content"].startswith("# Lesson Storyboard")


def test_extract_state_file_does_not_guess_unknown_markdown_document():
    original = "```markdown\n# Personal Notes\n\n- Keep this visible\n```"

    content, spec = _extract_state_file(original)

    assert content == original
    assert spec is None


def test_skill_validator_passes_on_the_skills_own_valid_fixture(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    _write_state_file(
        settings, "project-valid-state1", "alignment-map.md",
        _skill_fixture("alignment", "valid.md"),
    )

    report = _validate_project_state(settings, "project-valid-state1")

    assert report["status"] == "pass"
    assert report["checks"][0]["script"] == "scripts/validate_alignment_map.py"
    assert report["checks"][0]["exit_code"] == 0


def test_skill_validator_reports_findings_from_an_invalid_fixture(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    _write_state_file(
        settings, "project-invalid-stat", "alignment-map.md",
        _skill_fixture("alignment", "invalid.md"),
    )

    report = _validate_project_state(settings, "project-invalid-stat")

    assert report["status"] in {"fail", "incomplete"}
    findings = report["checks"][0]["findings"]
    assert findings, "the validator's stdout should be parsed into findings"
    assert {finding["level"] for finding in findings} & {"error", "gap"}


def test_validator_folds_lookalike_dashes_after_skill_sanitization(monkeypatch, tmp_path):
    """The updated skill folds typographic look-alikes instead of failing them, and the
    app's tool contract now coaches ASCII hyphens for identifiers only. Pin the
    relaxation so a future skill change that re-tightens it surfaces here."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    brief = _skill_fixture("design_state", "valid.md").replace("Co-design", "Co–design")
    assert "Co–design" in brief  # en dash where the schema expects a hyphen
    _write_state_file(settings, "project-endash-brief", "course-design-brief.md", brief)

    outcome = server._validate_state_file(
        settings, "project-endash-brief", "course-design-brief.md", "establish"
    )

    assert outcome["status"] == "pass"


def test_skill_tool_contract_matches_validator_separator_rules():
    """The prompt coaches the exact corrections the validators still enforce."""
    prompt, _ = _load_skill_runtime(get_settings(), "design")

    assert "semicolons, never commas" in prompt
    assert "plain ASCII hyphen" in prompt
    # The heading-dash coaching was retired when the skill started folding look-alikes.
    assert "column headings" not in server.SKILL_TOOL_CONTRACT


def test_curriculum_map_validation_includes_practice_distribution(monkeypatch, tmp_path):
    """The skill's curriculum-map template says --check-practice-distribution automates
    one of its review rows; the app is the only validator runner in its workflow, so
    the app must pass the flag or the automated check never runs anywhere."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    valid = _skill_fixture("course_curriculum_map", "valid.md")
    _write_state_file(settings, "project-practice-dist", "course-curriculum-map.md", valid)

    outcome = server._validate_state_file(
        settings, "project-practice-dist", "course-curriculum-map.md"
    )
    assert outcome["status"] == "pass"

    # Same outcome practiced twice in the same week, never distributed.
    massed = valid + (
        "| 4 | 2 | LO-1 | practice | none | Drill set A | Instructor feedback | 1 | approved |\n"
        "| 5 | 2 | LO-1 | practice | none | Drill set B | Instructor feedback | 1 | approved |\n"
    )
    _write_state_file(settings, "project-practice-dist", "course-curriculum-map.md", massed)

    outcome = server._validate_state_file(
        settings, "project-practice-dist", "course-curriculum-map.md"
    )
    assert outcome["status"] == "incomplete"
    assert any("massed" in finding["message"] for finding in outcome["findings"])


def test_skill_tools_are_offered_to_the_model():
    tools = server.skill_tool_definitions(get_settings())
    names = {tool["function"]["name"] for tool in tools}

    assert names == {
        "read_skill_template",
        "read_skill_reference",
        "validate_dataset",
        "list_state_files",
        "read_state_file",
        "write_state_file",
        "run_validators",
    }
    # The model must not be able to name a file the write allow-list would refuse.
    write = next(
        tool for tool in tools if tool["function"]["name"] == "write_state_file"
    )
    assert set(write["function"]["parameters"]["properties"]["file"]["enum"]) == set(
        STATE_FILE_VALIDATORS
    )
    # The reference reader offers exactly the installed references, so every
    # "Read references/X" instruction in SKILL.md is followable.
    reader = next(
        tool for tool in tools if tool["function"]["name"] == "read_skill_reference"
    )
    installed = {
        path.name
        for path in (get_settings().skill_dir / "references").glob("*.md")
    }
    assert set(reader["function"]["parameters"]["properties"]["file"]["enum"]) == installed


def test_write_state_file_tool_returns_validator_findings(monkeypatch, tmp_path):
    """The agent has to see why a write failed, or it cannot correct itself."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()

    result = server._dispatch_skill_tool(
        settings,
        "project-tooldispatch1",
        "write_state_file",
        {"file": "alignment-map.md", "content": _skill_fixture("alignment", "invalid.md")},
    )

    assert result["validation"]["status"] in {"fail", "incomplete"}
    assert result["validation"]["findings"], "findings must reach the model"

    good = server._dispatch_skill_tool(
        settings,
        "project-tooldispatch1",
        "write_state_file",
        {"file": "alignment-map.md", "content": _skill_fixture("alignment", "valid.md")},
    )
    assert good["validation"]["status"] == "pass"


def test_skill_tools_refuse_unknown_files_and_missing_projects(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()

    escaped = server._dispatch_skill_tool(
        settings, "project-tooldispatch2", "write_state_file",
        {"file": "../../escaped.md", "content": "x"},
    )
    assert "error" in escaped
    assert not (tmp_path / "escaped.md").exists()

    no_project = server._dispatch_skill_tool(
        settings, None, "run_validators", {"design_profile": "establish"}
    )
    assert "error" in no_project

    unknown = server._dispatch_skill_tool(
        settings, "project-tooldispatch2", "not_a_tool", {}
    )
    assert "error" in unknown


def test_read_skill_template_tool_serves_the_installed_asset():
    result = server._dispatch_skill_tool(
        get_settings(), None, "read_skill_template", {"file": "alignment-map.md"}
    )

    assert result["validator"] == "validate_alignment_map.py"
    assert "Cognitive demand" in result["template"]


def test_read_skill_reference_tool_serves_unrouted_references():
    """SKILL.md routes the model to references only some profiles pre-load; the
    reader keeps every such instruction followable instead of aspirational."""
    result = server._dispatch_skill_tool(
        get_settings(), None, "read_skill_reference",
        {"file": "rich-artifact-production.md"},
    )

    assert result["file"] == "references/rich-artifact-production.md"
    assert "production" in result["content"].lower()

    refused = server._dispatch_skill_tool(
        get_settings(), None, "read_skill_reference", {"file": "../SKILL.md"}
    )
    assert "error" in refused


def test_validate_dataset_tool_runs_the_skill_dataset_validator(monkeypatch, tmp_path):
    """SKILL.md requires the exact-dataset check before releasing a data-based
    activity; the tool runs it against the uploaded source."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    project_id = "project-datasettool1"
    source_id = "SRC-ABCDEF012345"
    source_dir = server._source_dir(settings, project_id, source_id)
    source_dir.mkdir(parents=True)
    # Sequential integers (1, 2, 3, ...) read as a row identifier to the
    # validator's data-quality heuristic, not a measurement -- use varied
    # floats so the check exercises paired-quantitative fit, not that gap.
    rows = "\n".join(f"{1.2 + i * 0.37:.2f},{4.1 * (1.2 + i * 0.37) - 0.1:.2f}" for i in range(12))
    (source_dir / "original.csv").write_text(f"x,y\n{rows}\n", encoding="utf-8")

    result = server._dispatch_skill_tool(
        settings, project_id, "validate_dataset",
        {"source_id": source_id, "representation": "scatter", "x": "x", "y": "y"},
    )
    assert result["script"] == "scripts/validate_dataset.py"
    assert result["status"] == "pass"
    # assets/data-task-record.md's "Dataset file" column wants a path relative
    # to the record itself, which lives in the project's state directory.
    assert result["dataset_path_relative_to_state_dir"] == os.path.join(
        "..", "sources", source_id, "original.csv"
    )

    missing = server._dispatch_skill_tool(
        settings, project_id, "validate_dataset",
        {"source_id": "SRC-000000000000", "representation": "scatter"},
    )
    assert "error" in missing


def test_validate_dataset_tool_passes_role_and_sheet_arguments(monkeypatch, tmp_path):
    """scripts/validate_dataset.py now pairs roles (--x/--y/...) rather than
    accepting bare --column for representations that require paired columns;
    the generic 'columns' list alone is no longer enough for a scatter plot."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    project_id = "project-datasettool2"
    source_id = "SRC-123456789ABC"
    source_dir = server._source_dir(settings, project_id, source_id)
    source_dir.mkdir(parents=True)
    (source_dir / "original.csv").write_text(
        "region,total\n" + "\n".join(f"r{i},{i * 3}" for i in range(1, 9)) + "\n",
        encoding="utf-8",
    )

    result = server._dispatch_skill_tool(
        settings, project_id, "validate_dataset",
        {
            "source_id": source_id,
            "representation": "bar",
            "category": "region",
            "value": "total",
        },
    )
    assert result["status"] == "pass"


def test_dataset_representation_enum_matches_the_installed_script_exactly():
    """A drift in either direction is a real bug: a token missing from the app
    means the model can never request it; a stale token means the app offers a
    representation the installed validator no longer accepts."""
    help_text = subprocess.run(
        [sys.executable, "scripts/validate_dataset.py", "--help"],
        cwd=str(get_settings().skill_dir),
        capture_output=True, text=True, check=True,
    ).stdout
    match = re.search(r"--representation\s+\{([^}]+)\}", help_text)
    installed = set(match.group(1).split(","))
    assert set(server.DATASET_REPRESENTATIONS) == installed


def test_dataset_column_roles_match_the_installed_script():
    help_text = subprocess.run(
        [sys.executable, "scripts/validate_dataset.py", "--help"],
        cwd=str(get_settings().skill_dir),
        capture_output=True, text=True, check=True,
    ).stdout
    for role in server.DATASET_COLUMN_ROLES:
        assert f"--{role} " in help_text, role


def test_handoff_state_files_now_run_their_validator(monkeypatch, tmp_path):
    """The skill update added validate_handoff_state.py for the design log, source
    register, and capability manifest; a write must return its findings."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    template = (settings.skill_dir / "assets" / "design-log.md").read_text(
        encoding="utf-8"
    )

    result = server._dispatch_skill_tool(
        settings, "project-handoffstate1", "write_state_file",
        {"file": "design-log.md", "content": template},
    )

    validation = result["validation"]
    assert validation["script"] == "scripts/validate_handoff_state.py"
    assert validation["status"] == "incomplete"
    assert any(
        "decision" in finding["message"].lower()
        for finding in validation["findings"]
    )


def test_data_task_record_is_registered_and_routed(monkeypatch, tmp_path):
    """The skill's re-executable data-task-fit chain: a manifest row claiming
    the data-task-fit token must link a row in data-task-record.md, and
    scripts/validate_data_task_record.py re-runs it against the dataset."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    assert STATE_FILE_VALIDATORS["data-task-record.md"] == "validate_data_task_record.py"

    template = server._dispatch_skill_tool(
        settings, None, "read_skill_template", {"file": "data-task-record.md"}
    )
    assert template["validator"] == "validate_data_task_record.py"

    result = server._dispatch_skill_tool(
        settings, "project-datatask1", "write_state_file",
        {"file": "data-task-record.md", "content": template["template"]},
    )
    assert result["validation"]["script"] == "scripts/validate_data_task_record.py"

    _, runtime = _load_skill_runtime(settings, "data")
    assert "assets/data-task-record.md" in runtime["loaded_assets"]


def test_data_task_record_confines_paths_to_the_true_project_root(monkeypatch, tmp_path):
    """scripts/validate_data_task_record.py checks that a recorded dataset stays
    inside the project, but defaults to treating the record's own directory
    (state/) as that boundary. This app stores uploaded datasets in a sibling
    sources/ directory, so a correctly-authored row referencing
    ../sources/<SRC-ID>/original.csv would be wrongly rejected as an escape
    unless the app supplies the true project root via --root."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    project_id = "project-datataskroot1"
    source_id = "SRC-CCCCCCCCCCCC"
    source_dir = server._source_dir(settings, project_id, source_id)
    source_dir.mkdir(parents=True)
    rows = "\n".join(f"{1.2 + i * 0.37:.2f},{4.1 * (1.2 + i * 0.37) - 0.1:.2f}" for i in range(12))
    csv_path = source_dir / "original.csv"
    csv_path.write_text(f"mass_kg,extension_mm\n{rows}\n", encoding="utf-8")
    digest = hashlib.sha256(csv_path.read_bytes()).hexdigest()

    content = (
        "# Data–Task Fit Record\n\n- Schema version: 1.0\n- Last updated: 2026-08-09\n\n"
        "| Artifact ID | Dataset file | Dataset SHA-256 | Dataset version or date | Worksheet "
        "| Representation | Column roles | Expected student output | Intended interpretation "
        "| Execution method | Execution evidence | Executed on | Result |\n"
        "|---|---|---|---|---|---|---|---|---|---|---|---|---|\n"
        f"| WS-1 | ../sources/{source_id}/original.csv | {digest} | 2026-08-09 |  | scatter "
        "| x=mass_kg; y=extension_mm | Fitted line with slope | Extension rises with mass "
        "| manual |  | 2026-08-09 | Produced; slope about 4.1 |\n"
    )
    result = server._dispatch_skill_tool(
        settings, project_id, "write_state_file",
        {"file": "data-task-record.md", "content": content},
    )
    assert result["validation"]["status"] == "pass", result["validation"]["findings"]


def test_release_records_run_validate_release_record(monkeypatch, tmp_path):
    """safety-review.md, accessibility-review.md, and production-plan.md used to
    have no validator at all -- a copied, unfilled template could reach
    teaching-ready status. scripts/validate_release_record.py now checks each
    carries an affirmative, unblocked decision with a named owner and dates."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    for filename in ("safety-review.md", "accessibility-review.md", "production-plan.md"):
        assert STATE_FILE_VALIDATORS[filename] == "validate_release_record.py", filename

        template = server._dispatch_skill_tool(
            settings, None, "read_skill_template", {"file": filename}
        )
        assert template["validator"] == "validate_release_record.py"

        result = server._dispatch_skill_tool(
            settings, "project-releaserecord1", "write_state_file",
            {"file": filename, "content": template["template"]},
        )
        validation = result["validation"]
        assert validation["script"] == "scripts/validate_release_record.py", filename
        # A blank copied template must never look like a completed review.
        assert validation["status"] == "incomplete", filename
        assert validation["findings"], filename


def test_handoff_state_validation_goes_strict_only_at_handoff_profile(monkeypatch, tmp_path):
    """validate_project.py only passes --strict once the project has reached the
    handoff design profile; the app's per-file validation must match, or a
    source register that is fine mid-project would block earlier phases."""
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    project_id = "project-strictcheck1"
    # Filled enough to satisfy the core (non-strict) fields but missing the
    # handoff-only provenance columns (owner/publisher, stable reference, ...).
    content = (
        "# Source Register\n\n- Schema version: 1.0\n- Last updated: 2026-01-01\n\n"
        "| Source ID | Title | Authority type | Owner/publisher | Publication/revision date | "
        "Stable reference | Last verified | Population/context and fit | Strength/limits | "
        "License/reuse | Supported claim or artifact | Status |\n"
        "|---|---|---|---|---|---|---|---|---|---|---|---|\n"
        "| SRC-1 | A study | scholarly | | | | | | | | Backs the sequencing choice | verified |\n"
    )
    server._write_state_file(settings, project_id, "source-register.md", content)

    lenient = server._validate_state_file(
        settings, project_id, "source-register.md", design_profile="establish"
    )
    strict = server._validate_state_file(
        settings, project_id, "source-register.md", design_profile="handoff"
    )
    assert lenient["status"] == "pass"
    assert strict["status"] == "incomplete"
    assert any("owner/publisher" in f["message"] for f in strict["findings"])


def test_guided_rapid_auto_modes_load_the_interaction_protocol_everywhere():
    """The mode instruction tells the model to apply the interaction-protocol
    rules; outside the establish profile those rules must still reach the prompt."""
    settings = get_settings()
    for mode in ("Guided", "Rapid", "Auto"):
        _, runtime = _load_skill_runtime(settings, "artifact", mode)
        assert "references/interaction-protocol.md" in runtime["loaded_files"], mode
    _, co_design = _load_skill_runtime(settings, "artifact", "Co-design")
    assert "references/interaction-protocol.md" not in co_design["loaded_files"]


def test_auto_mode_supplies_the_design_log_template():
    """Auto must record every self-answered decision card in design-log.md, so
    the template rides along regardless of profile."""
    _, runtime = _load_skill_runtime(get_settings(), "artifact", "Auto")
    assert "assets/design-log.md" in runtime["loaded_assets"]


def test_auto_instruction_requires_recorded_cards_and_state():
    instruction = _collaboration_mode_instruction("Auto")
    assert "do not ask the educator a question" in instruction
    assert "design-log.md" in instruction
    assert "unpresented, not unwritten" in instruction
    # Tier-scoping: the skill moved from "always design-log.md" to "the tier
    # decides where the record lives" -- Focused tier folds it into the
    # deliverable instead of creating the state bundle.
    assert "Focused tier" in instruction
    assert "do not create the state bundle" in instruction


def test_profile_carries_forward_when_a_continuation_has_no_keywords():
    prior_rows = [
        {
            "role": "assistant",
            "metadata_json": json.dumps({"skill_runtime": {"profile": "artifact"}}),
        },
        {"role": "user", "metadata_json": "{}"},
    ]

    profile = server._resolve_skill_profile(
        "auto",
        [ChatMessage(role="user", content="Yes, produce it as we discussed.")],
        prior_rows,
    )
    assert profile == "artifact"

    # A topic change names its own keywords and re-routes normally.
    profile = server._resolve_skill_profile(
        "auto",
        [ChatMessage(role="user", content="Now draft the quiz rubric.")],
        prior_rows,
    )
    assert profile == "assessment"

    # No history: the continuation stays at establish.
    assert server._resolve_skill_profile(
        "auto", [ChatMessage(role="user", content="Yes, go ahead.")], []
    ) == "establish"


def test_dataset_requests_route_to_the_data_profile_with_the_fit_reference():
    """SKILL.md requires data-task-fit.md whenever students receive a dataset,
    a spreadsheet, or a chart to produce."""
    for text in (
        "Build a worksheet where students plot a scatter chart from this dataset",
        "Make a spreadsheet activity on unit conversions",
        "Students should produce a histogram of the survey results",
    ):
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == "data", text

    _, runtime = _load_skill_runtime(get_settings(), "data")
    assert "references/data-task-fit.md" in runtime["loaded_files"]
    assert "assets/source-register.md" in runtime["loaded_assets"]
    assert "assets/data-task-record.md" in runtime["loaded_assets"]


def test_description_scope_keywords_reach_their_profiles():
    """The skill's own description names work the router previously missed."""
    cases = {
        "Design a peer evaluation for the team project": "assessment",
        "Students submit a peer review of each lab report": "assessment",
        "Draft the solution key for problem set 3": "assessment",
        "My students lose motivation halfway through the term": "design",
        "How do I build belonging in my first-year section?": "design",
        "Storyboard the lecture video introduction": "artifact",
        "Plan the capstone milestones": "course",
        "Does this handout work with a screen reader?": "accessibility",
    }
    for text, expected in cases.items():
        assert _infer_skill_profile([ChatMessage(role="user", content=text)]) == expected, text


def test_state_file_round_trip_through_the_api(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_id = "project-state-api01"
    client.post("/api/projects", json={"name": "State API"})

    written = client.put(
        f"/api/projects/{project_id}/state/alignment-map.md",
        json={"file": "alignment-map.md", "content": _skill_fixture("alignment", "valid.md")},
    )

    assert written.status_code == 200
    assert written.json()["state_file"]["validation"]["status"] == "pass"

    listed = client.get(f"/api/projects/{project_id}/state").json()
    assert [entry["file"] for entry in listed["state_files"]] == ["alignment-map.md"]

    read_back = client.get(f"/api/projects/{project_id}/state/alignment-map.md")
    assert read_back.status_code == 200
    assert read_back.json()["content"] == _skill_fixture("alignment", "valid.md").strip() + "\n"

    validated = client.post(f"/api/projects/{project_id}/validate").json()
    assert validated["status"] == "pass"
    assert "does not replace" in validated["scope_note"]

    deleted = client.delete(f"/api/projects/{project_id}/state/alignment-map.md")
    assert deleted.status_code == 200
    assert client.get(f"/api/projects/{project_id}/state").json()["state_files"] == []


def test_state_path_helper_refuses_names_outside_the_allow_list(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()

    for hostile in ("../../escaped.md", "../project-index.md", "notes.md", "/etc/passwd"):
        with pytest.raises(HTTPException) as raised:
            server._state_file_path(settings, "project-traversal01", hostile)
        assert raised.value.status_code == 400


def test_state_api_rejects_path_traversal(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    escape_target = tmp_path / "projects" / "escaped.md"

    response = client.put(
        "/api/projects/project-traversal01/state/..%2F..%2Fescaped.md",
        json={"file": "../../escaped.md", "content": "x"},
    )

    assert response.status_code in {400, 404}
    assert not escape_target.exists()
    assert not (tmp_path / "escaped.md").exists()


def test_skill_asset_endpoint_serves_the_installed_template():
    listed = client.get("/api/skill/assets")
    assert listed.status_code == 200
    files = {asset["file"] for asset in listed.json()["assets"]}
    assert "alignment-map.md" in files

    asset = client.get("/api/skill/assets/alignment-map.md")
    assert asset.status_code == 200
    assert asset.json()["validator"] == "validate_alignment_map.py"
    assert "Schema version" in asset.json()["content"]

    assert client.get("/api/skill/assets/not-a-real-asset.md").status_code == 404


def test_skill_runtime_adds_limited_model_reliability_checks():
    prompt, runtime = _load_skill_runtime(get_settings(), "artifact")

    assert "Reliability overlay for this limited-capability model" in prompt
    assert "displayed block totals equal the stated available time" in prompt
    assert "Cross-check every disciplinary example" in prompt
    assert runtime["reliability_overlay"]["target"] == "gpt-oss-120b"


def test_upload_extracts_and_lists_text_source(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_id = "testproject01"

    upload_response = client.post(
        f"/api/projects/{project_id}/sources",
        data={"data_classification_ack": "true"},
        files={
            "files": (
                "course-notes.txt",
                b"Students compare two reaction-rate models and justify a selection.",
                "text/plain",
            )
        },
    )

    assert upload_response.status_code == 200
    result = upload_response.json()
    assert result["errors"] == []
    assert result["sources"][0]["source_id"].startswith("SRC-")
    assert result["sources"][0]["filename"] == "course-notes.txt"
    assert result["sources"][0]["vector_index"]["algorithm"] == "local-hashed-term-cosine-v1"
    assert result["sources"][0]["vector_index"]["dimensions"] == 384
    vector_path = (
        tmp_path
        / "projects"
        / project_id
        / "sources"
        / result["sources"][0]["source_id"]
        / "vectors.json"
    )
    assert vector_path.is_file()

    list_response = client.get(f"/api/projects/{project_id}/sources")
    assert list_response.status_code == 200
    assert len(list_response.json()["sources"]) == 1

    source_id = result["sources"][0]["source_id"]
    preview_response = client.get(f"/api/projects/{project_id}/sources/{source_id}/preview")
    assert preview_response.status_code == 200
    assert "reaction-rate models" in preview_response.json()["preview_text"]
    assert preview_response.json()["vector_chunks"][0]["main_idea"]
    assert "reaction-rate" in preview_response.json()["vector_chunks"][0]["keywords"]
    assert preview_response.json()["vector_metadata"]["metadata_algorithm"] == "extractive-main-idea-v3"

    download_response = client.get(f"/api/projects/{project_id}/sources/{source_id}/download")
    assert download_response.status_code == 200
    assert download_response.content == b"Students compare two reaction-rate models and justify a selection."
    assert "course-notes.txt" in download_response.headers["content-disposition"]


def test_upload_requires_data_confirmation(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))

    response = client.post(
        "/api/projects/testproject02/sources",
        data={"data_classification_ack": "false"},
        files={"files": ("notes.txt", b"Public course notes", "text/plain")},
    )

    assert response.status_code == 400
    assert "FERPA" in response.json()["detail"]


def test_chat_requires_api_key(monkeypatch):
    monkeypatch.setenv("PURDUE_GENAI_API_KEY", "")
    monkeypatch.setenv("PURDUE_GENAI_MODEL_ID", "gpt-oss-120b")

    response = client.post(
        "/api/chat",
        json={"messages": [{"role": "user", "content": "Hello"}]},
    )

    assert response.status_code == 503
    assert "PURDUE_GENAI_API_KEY" in response.json()["detail"]


def test_port_fallback_skips_occupied_port():
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as occupied:
        occupied.bind(("127.0.0.1", 0))
        occupied.listen(1)
        occupied_port = occupied.getsockname()[1]

        selected = find_available_port("127.0.0.1", occupied_port, 2)

    assert selected == occupied_port + 1


def test_project_workspace_persists_context(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    create_response = client.post(
        "/api/projects",
        json={
            "name": "CHE 205 redesign",
            "course_name": "Reaction engineering",
            "outcome": "Analyze rate data.",
            "mode": "Guided",
        },
    )

    assert create_response.status_code == 201
    project_id = create_response.json()["project"]["id"]

    update_response = client.patch(
        f"/api/projects/{project_id}",
        json={"class_time": "75 minutes", "notes": "Teams of four"},
    )
    workspace_response = client.get(f"/api/projects/{project_id}")

    assert update_response.status_code == 200
    assert workspace_response.status_code == 200
    workspace = workspace_response.json()
    assert workspace["project"]["name"] == "CHE 205 redesign"
    assert workspace["project"]["class_time"] == "75 minutes"
    assert workspace["project"]["notes"] == "Teams of four"
    assert workspace["messages"] == []
    assert workspace["artifacts"] == []


def test_markdown_renderer_supports_tables_and_sanitizes_html():
    rendered = _render_markdown(
        "# Artifact\n\n| Outcome | Evidence |\n|---|---|\n| Analyze | Explanation |"
        "\n\n<script>alert('unsafe')</script>"
    )

    assert "<h1>Artifact</h1>" in rendered
    assert "<table>" in rendered
    assert "<script>" not in rendered


def test_frontend_applies_responsive_semantic_table_treatment():
    app_js = (server.APP_DIR / "static" / "app.js").read_text(encoding="utf-8")
    styles = (server.APP_DIR / "static" / "styles.css").read_text(encoding="utf-8")

    assert "function enhanceMarkdownTables" in app_js
    assert 'table.classList.add("wide-table")' in app_js
    assert "cell.dataset.label" in app_js
    assert "table.wide-table" in styles
    assert "td::before" in styles
    assert "min-width: 130px" not in styles


def test_structured_decision_is_removed_and_parsed():
    content, decision = _extract_decision(
        "Choose an evidence pattern.\n\n```decision\n"
        '{"question":"Which pattern?","options":['
        '{"label":"Performance (Recommended)","description":"Authentic evidence",'
        '"value":"Use a performance task"},'
        '{"label":"Written explanation","description":"Fast to review",'
        '"value":"Use a written explanation"}]}\n```'
    )

    assert content == "Choose an evidence pattern."
    assert decision["question"] == "Which pattern?"
    assert len(decision["options"]) == 2


def test_bundled_decision_cards_are_all_removed_not_just_the_selected_one():
    """The skill's own eval added check_single_decision_per_turn to catch gpt-oss
    bundling several decision cards in one reply. The UI can only act on one, but
    leaving the others as raw JSON fences would show the professor unparsed
    leftovers instead of a clean message."""
    content, decision = _extract_decision(
        "Decision 1: pick a pattern.\n\n```decision\n"
        '{"question":"Which pattern?","options":['
        '{"label":"A (Recommended)","description":"First","value":"a"},'
        '{"label":"B","description":"Second","value":"b"}]}\n```\n\n'
        "Decision 2: pick a format.\n\n```decision\n"
        '{"question":"Which format?","options":['
        '{"label":"C (Recommended)","description":"Third","value":"c"},'
        '{"label":"D","description":"Fourth","value":"d"}]}\n```'
    )

    assert "```" not in content
    assert "question" not in content.lower() or "Which" not in content
    assert decision["question"] == "Which format?"  # last block wins, as before
    assert content == "Decision 1: pick a pattern.\n\nDecision 2: pick a format."


def test_suggested_replies_line_is_stripped_when_the_native_card_renders():
    """The skill's portable fallback line ('Suggested replies: "1" | ...') is for
    clients without a native choice control; next to the app's own decision buttons
    it is duplicate chrome."""
    content, decision = _extract_decision(
        "Pick a pattern.\n\n"
        'Suggested replies: "1" | "2" | "modify: <your change>" | "decide for me"\n\n'
        "```decision\n"
        '{"question":"Which pattern?","options":['
        '{"label":"A (Recommended)","description":"Best aligned","value":"a"},'
        '{"label":"B","description":"Alternative","value":"b"}]}\n```'
    )

    assert decision is not None
    assert content == "Pick a pattern."


def test_suggested_replies_line_survives_when_no_card_is_recovered():
    """Without a rendered card, typed replies are the educator's only cheap answer."""
    content, decision = _extract_decision(
        "Here is the plan summary.\n\n"
        'Suggested replies: "approve" | "modify: <your change>"'
    )

    assert decision is None
    assert "Suggested replies" in content


def test_interactive_modes_scope_approval_to_the_presented_piece():
    """One preview approval never authorizes the remaining artifact family; skipping
    checkpoints is a mode change only the educator can name."""
    for mode in ("Co-design", "Guided"):
        instruction = _collaboration_mode_instruction(mode)
        assert "only the piece just presented" in instruction, mode
        assert "mode change" in instruction, mode
    assert "piece just presented" not in _collaboration_mode_instruction("Auto")
    assert "piece just presented" not in _collaboration_mode_instruction("Rapid")


def test_decision_dialog_offers_decide_for_me_delegation():
    """The skill requires a native choice affordance to carry 'decide for me' as an
    explicit option, delegating the current decision to the recommendation."""
    app_js = (server.APP_DIR / "static" / "app.js").read_text(encoding="utf-8")
    index_html = (server.APP_DIR / "static" / "index.html").read_text(encoding="utf-8")

    assert 'id="delegateDecisionButton"' in index_html
    assert "delegateDecisionButton" in app_js
    assert "recommendedOption" in app_js


def test_new_project_defaults_to_co_design_mode(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    response = client.post("/api/projects", json={"name": "Mode default"})

    assert response.status_code == 201
    assert response.json()["project"]["mode"] == "Co-design"


def test_projects_saved_before_the_rename_migrate_to_co_design(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    settings = get_settings()
    with server._database_connection(settings) as connection:
        connection.execute(
            """
            INSERT INTO projects (id, name, mode, created_at, updated_at)
            VALUES ('project-legacy-mode', 'Pre-rename', 'Studio', '2026-01-01', '2026-01-01')
            """
        )

    response = client.get("/api/projects/project-legacy-mode")

    assert response.status_code == 200
    assert response.json()["project"]["mode"] == "Co-design"


def test_auto_mode_accepts_project_setting_and_resolves_choice_without_modal(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_response = client.post("/api/projects", json={"name": "Auto project", "mode": "Auto"})
    assert project_response.status_code == 201
    project = project_response.json()["project"]
    assert project["mode"] == "Auto"
    assert "do not ask the educator a question" in _collaboration_mode_instruction("Auto")

    content, interactive_decision = _extract_decision(
        "Select a path.\n\n```decision\n"
        '{"question":"Which path?","options":['
        '{"label":"A (Recommended)","description":"Best aligned","value":"a"},'
        '{"label":"B","description":"Alternative","value":"b"}]}'
        "\n```"
    )
    resolved, remaining_decision, auto_decision = _apply_auto_decision(content, interactive_decision)
    assert remaining_decision is None
    assert auto_decision["selected_label"] == "A (Recommended)"
    assert "Auto decision recorded" in resolved
    assert "Best aligned" in resolved


def test_question_list_falls_back_to_choice_dialog():
    content, decision = _extract_decision(
        "Please choose:\n"
        "1. **Which strategies should students compare?** "
        "(e.g., case-based learning vs. lecture, problem-based learning vs. flipped classroom)"
    )

    assert content.startswith("Please choose")
    assert decision["question"] == "Which strategies should students compare?"
    assert decision["options"][0]["label"].endswith("(Recommended)")
    assert len(decision["options"]) == 3


def test_generic_json_fence_is_accepted_for_decision_schema():
    content, decision = _extract_decision(
        "Choose a direction.\n```json\n"
        '{"question":"Which direction?","options":['
        '{"label":"A (Recommended)","description":"Use A","value":"A"},'
        '{"label":"B","description":"Use B","value":"B"}]}\n```'
    )

    assert content == "Choose a direction."
    assert decision["question"] == "Which direction?"


def test_trailing_bare_decision_json_is_recovered_and_hidden():
    content, decision = _extract_decision(
        "The lesson plan is ready.\n\n"
        '{"question":"Which feedback method?","options":['
        '{"label":"Cards (Recommended)","description":"Fast visual check","value":"cards"},'
        '{"label":"Poll","description":"Saved responses","value":"poll"}]}'
    )

    assert content == "The lesson plan is ready."
    assert decision["question"] == "Which feedback method?"


def test_explicit_artifact_request_wins_over_review_keyword():
    profile = _infer_skill_profile(
        [ChatMessage(role="user", content="Create an artifact with an instructor review note")]
    )

    assert profile == "artifact"


def test_open_decision_is_not_saved_as_finished_artifact(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_id = "project-decision-test"

    _, assistant, artifact = _persist_exchange(
        get_settings(),
        project_id,
        "Create a worksheet artifact",
        "Which worksheet pattern should we use?",
        {
            "decision": {
                "question": "Which pattern?",
                "options": [
                    {"label": "A", "description": "A", "value": "A"},
                    {"label": "B", "description": "B", "value": "B"},
                ],
            }
        },
        "artifact",
    )

    assert assistant["decision"] is not None
    assert artifact is None
    assert client.get(f"/api/projects/{project_id}").json()["artifacts"] == []


def test_specific_message_download_includes_skill_route(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_id = "project-message-download"
    _, assistant, _ = _persist_exchange(
        get_settings(),
        project_id,
        "Draft a table",
        "## Comparison\n\n| A | B |\n|---|---|\n| 1 | 2 |",
        {
            "skill_runtime": {
                "profile": "artifact",
                "loaded_files": ["SKILL.md", "references/artifact-patterns.md"],
                "fingerprint": "testfingerprint01",
            },
            "sources_used": [],
        },
        "artifact",
    )

    response = client.get(
        f"/api/projects/{project_id}/messages/{assistant['id']}/download?format=markdown"
    )

    assert response.status_code == 200
    assert "SKILL profile: artifact" in response.text
    assert "references/artifact-patterns.md" in response.text
    assert "## Comparison" in response.text


def test_skill_trace_links_decision_to_selected_answer(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_id = "project-trace-test"
    _, decision_message, _ = _persist_exchange(
        get_settings(),
        project_id,
        "Help me choose",
        "Choose an evidence pattern.",
        {
            "decision": {
                "question": "Which evidence pattern?",
                "options": [
                    {"label": "Performance", "description": "Authentic", "value": "performance"},
                    {"label": "Explanation", "description": "Written", "value": "explanation"},
                ],
            },
            "skill_runtime": {
                "profile": "design",
                "loaded_files": ["SKILL.md", "references/design-workflow.md"],
                "fingerprint": "tracefingerprint1",
            },
        },
        "design",
    )
    _persist_exchange(
        get_settings(),
        project_id,
        "Use a performance task",
        "The performance task is now the working evidence pattern.",
        {
            "skill_runtime": {
                "profile": "design",
                "loaded_files": ["SKILL.md", "references/design-workflow.md"],
                "fingerprint": "tracefingerprint1",
            }
        },
        "design",
        "Selected: Performance",
        {
            "origin_message_id": decision_message["id"],
            "question": "Which evidence pattern?",
            "selected_label": "Performance",
            "selected_value": "performance",
        },
    )

    response = client.get(f"/api/projects/{project_id}/trace")

    assert response.status_code == 200
    trace = response.json()
    assert trace["summary"]["decisions"] == 1
    assert trace["summary"]["answers"] == 1
    answer_event = next(event for event in trace["events"] if event["type"] == "answer")
    assert answer_event["decision_trace"]["origin_message_id"] == decision_message["id"]
    assert answer_event["decision_trace"]["selected_label"] == "Performance"


def test_project_export_downloads_markdown(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project = client.post("/api/projects", json={"name": "Export test"}).json()["project"]

    response = client.get(f"/api/projects/{project['id']}/export?format=markdown")

    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/markdown")
    assert "attachment" in response.headers["content-disposition"]
    assert "# Export test" in response.text


def test_artifact_spec_contract_is_extracted_and_hidden_from_chat():
    content = """Your worksheet is ready.

```artifact_spec
{"kind":"worksheet","title":"Evidence Check","subtitle":"Practice","sections":[{"heading":"Claim","body":"","bullets":[],"prompts":["What is the strongest evidence?"],"checklist":[],"response_lines":3,"table":null}],"source_ids":[]}
```"""

    cleaned, spec = _extract_artifact_spec(content)

    assert cleaned == "Your worksheet is ready."
    assert spec["kind"] == "worksheet"
    assert spec["sections"][0]["prompts"] == ["What is the strongest evidence?"]


def test_generic_json_artifact_contract_is_recovered_and_hidden():
    content = """```json
{"kind":"worksheet","title":"Loop Check","sections":[{"heading":"Question","body":"Explain the invariant."}],"source_ids":[]}
```"""

    cleaned, spec = _extract_artifact_spec(content)

    assert spec["kind"] == "worksheet"
    assert spec["title"] == "Loop Check"
    assert "Loop Check" in cleaned and "artifact generation" not in cleaned
    assert "```json" not in cleaned


def test_bare_artifact_contract_accepts_instructor_key_without_response_space():
    bullets = [f"Answer {index}" for index in range(10)]
    payload = json.dumps(
        {
            "kind": "worksheet",
            "title": "Loop Quiz",
            "sections": [
                {
                    "heading": "Instructor answer key",
                    "bullets": bullets,
                    "response_lines": 0,
                }
            ],
            "source_ids": [],
        }
    )

    cleaned, spec = _extract_artifact_spec(payload)

    assert spec["title"] == "Loop Quiz"
    assert spec["sections"][0]["response_lines"] == 0
    assert spec["sections"][0]["bullets"] == bullets
    assert "Loop Quiz" in cleaned and "artifact generation" not in cleaned


def test_incomplete_internal_json_is_replaced_with_plain_recovery_status():
    content = '```json\n{"kind":"worksheet","sections":[]}\n```'

    cleaned = _hide_incomplete_internal_payloads(content)

    assert "internal production payload" in cleaned
    assert "```json" not in cleaned
    assert '"sections"' not in cleaned


def test_incomplete_bare_internal_json_is_not_professor_visible():
    content = '{"kind":"worksheet","sections":[]}'

    cleaned = _hide_incomplete_internal_payloads(content)

    assert "internal production payload" in cleaned
    assert '"sections"' not in cleaned


def test_markdown_document_fence_is_unwrapped_without_touching_python_code():
    content = (
        "Here is the lesson.\n\n```markdown\n# Loops\n\n"
        "| Goal | Evidence |\n|---|---|\n| Trace a loop | Explanation |\n```"
        "\n\n```python\nfor item in items:\n    print(item)\n```"
    )

    cleaned = _unwrap_markdown_documents(content)
    rendered = _render_markdown(content)

    assert "```markdown" not in cleaned
    assert "```python" in cleaned
    assert "<h1>Loops</h1>" in rendered
    assert "<table>" in rendered
    assert "<pre><code>" in rendered


def test_legacy_professor_content_is_normalized_before_display():
    content = "```markdown\n# Quiz review\n\n| Item | Note |\n|---|---|\n| 1 | Revise |\n```"

    display = _prepare_professor_display_content(content)

    assert display.startswith("# Quiz review")
    assert "```markdown" not in display


def test_assessment_decision_never_asks_professor_for_internal_metadata():
    content, decision = _sanitize_professor_decision(
        "What information can you provide?",
        {
            "question": "Which production route?",
            "options": [
                {
                    "label": "Provide outcome ID and cognitive demand",
                    "description": "Creates an assessment blueprint",
                    "value": "metadata",
                },
                {"label": "Provide schema", "description": "JSON", "value": "schema"},
            ],
        },
        "assessment",
    )

    assert content == "One teaching choice will help me draft the quiz appropriately."
    assert decision["question"] == "How will students use this quiz?"
    assert [option["label"] for option in decision["options"]] == [
        "Practice and feedback (Recommended)",
        "Counts toward the grade",
    ]
    assert "outcome" not in str(decision).lower()


def test_redundant_assessment_decision_is_dropped_after_three_plain_questions():
    content, decision = _sanitize_professor_decision(
        "Is this practice or graded?\nWhat should students explain?\nHow many minutes?",
        {
            "question": "Select the key quiz parameters:",
            "options": [
                {"label": "Practice", "description": "Feedback", "value": "practice"},
                {"label": "Graded", "description": "Scores", "value": "graded"},
            ],
        },
        "assessment",
    )

    assert content.count("?") == 3
    assert decision is None


def test_duplicate_lesson_question_is_not_repeated_as_decision_card():
    question = (
        "Assessment goal – Should students demonstrate mastery through a problem set or a quiz?"
    )
    content, decision = _sanitize_professor_decision(
        "1. What do students know?\n"
        "2. How much class time is available?\n"
        f"3. **{question}**",
        {
            "question": question,
            "options": [
                {"label": "Quiz", "description": "Short check", "value": "quiz"},
                {"label": "Problem set", "description": "More practice", "value": "problems"},
            ],
        },
        "establish",
    )

    assert content.count("?") == 3
    assert decision is None


def test_generic_lesson_confirmation_card_is_dropped_after_three_questions():
    content, decision = _sanitize_professor_decision(
        "What do students know?\nHow much time is available?\nWhich artifact is useful?",
        {
            "question": "Which details should I confirm to design the teaching material?",
            "options": [
                {"label": "All details", "description": "Complete", "value": "all"},
                {"label": "Defaults", "description": "Assume", "value": "defaults"},
            ],
        },
        "establish",
    )

    assert content.count("?") == 3
    assert decision is None


def test_overlong_clarification_is_reduced_to_three_professor_questions():
    response = _enforce_professor_question_limit(
        "Open questions that will shape the design:\n"
        "1. What language? What level?\n"
        "2. What should students do? What evidence?\n"
        "3. How much time?",
        "establish",
    )

    assert response.count("?") == 3
    assert "What should students be able to do" in response
    assert "class setting" in response


def test_natural_overlong_professor_questions_are_reduced_to_three():
    response = _enforce_professor_question_limit(
        "Here are a few quick questions to shape the materials—your answers will let me draft.\n"
        "Is this one lecture? How many minutes? What do students know? "
        "Would you like slides? Do you need a quiz?",
        "establish",
    )

    assert response.count("?") == 3
    assert "What should students be able to do" in response


def test_help_shape_interview_is_reduced_to_three_questions():
    response = _enforce_professor_question_limit(
        "These questions will help shape the lesson. What should students do? "
        "Which task matters? What evidence will demonstrate learning? Which assessment? "
        "What constraints shape the lesson? Which resources are available?",
        "establish",
    )

    assert response.count("?") == 3


def test_incomplete_assessment_payload_becomes_professor_clarification():
    response = _recover_professor_clarification(
        "> The app withheld an incomplete internal production payload. "
        "Regenerate this item to create a readable artifact.",
        "assessment",
    )

    assert response.count("?") == 3
    assert "practice and feedback" in response
    assert "internal" not in response.lower()


def test_artifact_preview_omits_empty_placeholder_table():
    preview = _artifact_spec_markdown(
        {
            "kind": "document",
            "title": "Lesson",
            "subtitle": "",
            "sections": [
                {
                    "heading": "Practice",
                    "body": "Try the activity.",
                    "bullets": [],
                    "prompts": [],
                    "checklist": [],
                    "response_lines": 0,
                    "table": {"headers": [""], "rows": [[""]]},
                }
            ],
            "source_ids": [],
        },
        {"course_name": "VLSI", "name": "VLSI lesson"},
    )

    assert "|  |" not in preview
    assert "Try the activity." in preview


def test_artifact_spec_repairs_model_latex_backslashes_and_allows_twelve_prompts():
    raw = r'''{
      "kind": "document",
      "title": "VLSI quiz",
      "subtitle": "Normalized practice",
      "sections": [{
        "heading": "Sizing comparison",
        "body": "Use \(d = g\,h + p\) and explain the result.",
        "prompts": ["1", "2", "3", "4", "5", "6", "7", "8", "9", "10", "11"],
        "table": {"headers": [], "rows": []}
      }],
      "source_ids": []
    }'''

    content, artifact_spec = _extract_artifact_spec(raw)

    # The reply summarizes the file rather than announcing that generation happened.
    assert "**VLSI quiz**" in content
    assert "artifact generation" not in content
    assert len(content) > 60
    assert artifact_spec["sections"][0]["body"] == r"Use \(d = g\,h + p\) and explain the result."
    assert len(artifact_spec["sections"][0]["prompts"]) == 11


def test_valid_artifact_spec_creates_artifact_even_on_general_profile(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project_id = "project-general-artifact"
    artifact_spec = {
        "kind": "worksheet",
        "title": "Revised loop quiz",
        "subtitle": "",
        "sections": [{"heading": "Question", "body": "Explain the invariant."}],
        "source_ids": [],
    }

    _, _, artifact = _persist_exchange(
        get_settings(),
        project_id,
        "Revise Question 2",
        "Prepared the revised quiz.",
        {"artifact_spec": artifact_spec},
        "establish",
    )

    assert artifact is not None
    assert artifact["message_id"]


def test_local_artifact_tools_create_slides_document_and_worksheet(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project = client.post(
        "/api/projects",
        json={"name": "Office tool test", "course_name": "EDCI  Test Studio"},
    ).json()["project"]
    expected_formats = {"slides": "pptx", "document": "docx", "worksheet": "docx"}

    for kind, file_format in expected_formats.items():
        response = client.post(
            f"/api/projects/{project['id']}/artifact-tools/generate",
            json={
                "kind": kind,
                "title": f"{kind.title()} smoke test",
                "subtitle": "Deterministic artifact harness",
                "sections": [
                    {
                        "heading": "Compare the evidence",
                        "body": "Use the supplied criteria to make an instructor-reviewed choice.",
                        "bullets": ["Identify the claim", "Connect evidence to the outcome"],
                        "prompts": ["Which evidence is most convincing?"] if kind == "worksheet" else [],
                        "checklist": ["Evidence is source-grounded"],
                        "response_lines": 2,
                        "table": {
                            "headers": ["Criterion", "Evidence"],
                            "rows": [["Alignment", "Directly supports the outcome"]],
                        },
                    }
                ],
                "source_ids": [],
            },
        )

        assert response.status_code == 200, response.text
        artifact = response.json()["artifact"]
        assert artifact["has_file"] is True
        assert artifact["file_format"] == file_format
        assert artifact["tool_trace"]["schema_validated"] is True
        assert "Compare the evidence" in artifact["html"]
        assert "Directly supports the outcome" in artifact["html"]
        download = client.get(
            f"/api/projects/{project['id']}/artifacts/{artifact['id']}/download?format=office"
        )
        assert download.status_code == 200
        assert len(download.content) > 1_000
        if kind == "slides":
            deck = Presentation(BytesIO(download.content))
            assert len(deck.slides) == 2
            assert "Slides smoke test" in deck.slides[0].shapes[0].text
            assert "[Sources]" in deck.slides[1].notes_slide.notes_text_frame.text
        else:
            document = Document(BytesIO(download.content))
            document_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
            assert f"{kind.title()} smoke test" in document_text
            assert "Compare the evidence" in document_text
            if kind == "worksheet":
                assert "Which evidence is most convincing?" in document_text

    project_package = client.get(f"/api/projects/{project['id']}/export?format=zip")
    assert project_package.status_code == 200
    assert project_package.headers["content-type"].startswith("application/zip")
    with zipfile.ZipFile(BytesIO(project_package.content)) as archive:
        names = archive.namelist()
        assert "conversation.md" in names
        assert "project.json" in names
        assert any(name.endswith(".pptx") for name in names)
        assert sum(name.endswith(".docx") for name in names) == 2


def test_legacy_decision_json_is_hidden_when_workspace_reloads(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project = client.post("/api/projects", json={"name": "Legacy decision"}).json()["project"]
    _persist_exchange(
        get_settings(),
        project["id"],
        "Help me choose",
        "Choose a direction.\n\n```json\n{\"question\":\"Which direction?\",\"options\":[{\"label\":\"A\",\"description\":\"First\",\"value\":\"a\"},{\"label\":\"B\",\"description\":\"Second\",\"value\":\"b\"}]}\n```",
        {"skill_runtime": {"profile": "design", "loaded_files": ["SKILL.md"], "fingerprint": "legacy"}},
        "design",
    )

    workspace = client.get(f"/api/projects/{project['id']}").json()
    assistant = workspace["messages"][-1]
    assert assistant["content"] == "Choose a direction."
    assert assistant["decision"]["question"] == "Which direction?"
    assert "```json" not in assistant["html"]


def test_artifact_tool_rejects_invented_source_ids(monkeypatch, tmp_path):
    monkeypatch.setenv("APP_DATA_DIR", str(tmp_path))
    project = client.post("/api/projects", json={"name": "Source guard"}).json()["project"]

    response = client.post(
        f"/api/projects/{project['id']}/artifact-tools/generate",
        json={
            "kind": "document",
            "title": "Unsafe source test",
            "sections": [{"heading": "Claim", "body": "A claim"}],
            "source_ids": ["SRC-FFFFFFFFFFFF"],
        },
    )

    assert response.status_code == 400
    assert "unknown project sources" in response.json()["detail"]


def test_settings_masks_api_key_and_updates_selected_env(monkeypatch, tmp_path):
    env_paths = {
        "home": tmp_path / "home.env",
        "app": tmp_path / "app.env",
        "current": tmp_path / "current.env",
    }
    monkeypatch.setattr(server, "_environment_paths", lambda: env_paths)
    monkeypatch.setenv("PURDUE_GENAI_API_KEY", "must-not-be-returned")

    get_response = client.get("/api/settings")
    patch_response = client.patch(
        "/api/settings",
        json={
            "target": "current",
            "values": {
                "PURDUE_GENAI_MODEL_ID": "gpt-oss:120b",
                "APP_PORT": "8001",
            },
        },
    )

    assert get_response.status_code == 200
    assert "must-not-be-returned" not in get_response.text
    api_key = next(
        field for field in get_response.json()["fields"]
        if field["name"] == "PURDUE_GENAI_API_KEY"
    )
    assert api_key["configured"] is True
    assert api_key["value"] == ""
    assert patch_response.status_code == 200
    assert patch_response.json()["restart_required"] is True
    assert "PURDUE_GENAI_MODEL_ID='gpt-oss:120b'" in env_paths["current"].read_text()


def test_latex_is_rewritten_into_readable_notation():
    """The app renders Markdown, not maths, so LaTeX reaches the educator verbatim."""
    rendered = _render_markdown(
        r"compute the electrical effort (h = C_{\text{load}}/C_{\text{in}}) per stage"
    )

    assert "C_load/C_in" in rendered
    assert "\\text" not in rendered
    assert "{" not in rendered


def test_latex_delimiters_and_symbols_are_converted():
    for source, expected in (
        (r"Apply $\frac{C_{out}}{C_{in}}$ now.", "C_out/C_in"),
        (r"We use \(d = gh + p\).", "d = gh + p"),
        (r"\[ E = \frac{1}{2} C V^2 \]", "1/2 C V^2"),
        (r"gain \times 2 and \tau \le 5 ns", "×"),
        (r"stage effort \sqrt{F}", "√(F)"),
    ):
        assert expected in server._plain_math(source), source


def test_plain_math_never_touches_code():
    """A Python lesson's code must survive byte for byte, backslashes included."""
    source = 'Trace it:\n\n```python\nprint("a\\tb")\ntotal += n\n```\n\nthen discuss.'

    assert server._plain_math(source) == source
    assert server._plain_math("Use `C_{\\text{in}}` here.") == "Use `C_{\\text{in}}` here."


def test_plain_math_leaves_ordinary_prose_alone():
    prose = "No mathematics here, just teaching prose about loops."

    assert server._plain_math(prose) == prose


def test_artifact_sections_without_content_are_dropped():
    """A dry run shipped a quiz with 'Questions' and 'Answer Key' and no questions."""
    spec = server.ArtifactToolRequest.model_validate(
        {
            "kind": "document",
            "title": "Formative Loop-Tracing Quiz",
            "sections": [
                {"heading": "Instructions", "body": "Work silently; show all work."},
                {"heading": "Questions"},
                {"heading": "Answer Key (Instructor Use Only)"},
            ],
        }
    )

    assert [section.heading for section in spec.sections] == ["Instructions"]


def test_artifact_with_no_renderable_section_is_rejected():
    with pytest.raises(ValueError):
        server.ArtifactToolRequest.model_validate(
            {
                "kind": "document",
                "title": "Hollow",
                "sections": [{"heading": "Questions"}, {"heading": "Answer Key"}],
            }
        )


def test_response_lines_alone_do_not_make_a_section_renderable():
    """Blank ruled space under a bare heading is an empty quiz, not a question."""
    with pytest.raises(ValueError):
        server.ArtifactToolRequest.model_validate(
            {
                "kind": "worksheet",
                "title": "Blank",
                "sections": [{"heading": "Question 1", "response_lines": 6}],
            }
        )


def test_a_table_only_section_is_kept():
    """Worksheet content often lives entirely in a table; that is real content."""
    spec = server.ArtifactToolRequest.model_validate(
        {
            "kind": "worksheet",
            "title": "Trace",
            "sections": [
                {
                    "heading": "Q1",
                    "table": {"headers": ["pass", "total"], "rows": [["1", "2"]]},
                }
            ],
        }
    )

    assert len(spec.sections) == 1


def test_internal_contract_names_never_reach_the_professor():
    """Prompt instructions alone did not hold, so the scrub is deterministic."""
    for source in (
        "| **Word file** (editable) – see `artifact_spec` below | notes |",
        "I will emit the artifact_spec now.",
        "Saved to state_file for you.",
    ):
        cleaned = server._strip_internal_references(source)
        assert "artifact_spec" not in cleaned
        assert "state_file" not in cleaned

    prose = "A normal sentence about worksheets and specifications."
    assert server._strip_internal_references(prose) == prose


def _realistic_arc():
    import importlib.util
    import sys

    if "realistic_arc" in sys.modules:
        return sys.modules["realistic_arc"]
    spec = importlib.util.spec_from_file_location(
        "realistic_arc", Path(__file__).resolve().parents[1] / "evaluations" / "realistic_arc.py"
    )
    module = importlib.util.module_from_spec(spec)
    # Register before exec: the module uses postponed annotations, and dataclasses
    # resolves them through sys.modules[cls.__module__] at class-creation time.
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def test_question_counter_ignores_drafted_student_material():
    """A reply that drafts quiz items is not interrogating the professor."""
    arc = _realistic_arc()

    asks_only = "A few things:\n1. What should students do?\n2. How long is the class?"
    with_draft = (
        "Two questions:\n1. How long is the class?\n2. Graded or practice?\n\n"
        "**Draft quiz**\n1. What is i on pass 3?\n2. What prints for x?"
    )

    assert arc.count_questions(asks_only) == 2
    assert arc.count_questions(with_draft) == 2  # the drafted items are not asks
    assert arc.count_questions("Here is the plan with timings.") == 0


def test_question_counter_does_not_count_parenthetical_examples():
    arc = _realistic_arc()

    assert arc.count_questions('How long is the lab? (e.g. "50 minutes?")') == 1
