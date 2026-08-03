from __future__ import annotations

from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptxInches, Pt as PptxPt


PURDUE_BLACK = "000000"
PURDUE_GOLD = "CFB991"
PURDUE_DARK_GOLD = "8E6F3E"
PURDUE_AGED = "F4EDDF"
INK = "202124"
MUTED = "5B5F63"
LIGHT_LINE = "D9D4CC"


def _set_run_font(run, *, name: str = "Arial", size: float = 11, bold: bool = False, color: str = INK) -> None:
    run.font.name = name
    run._element.get_or_add_rPr().rFonts.set(qn("w:ascii"), name)
    run._element.get_or_add_rPr().rFonts.set(qn("w:hAnsi"), name)
    run.font.size = Pt(size)
    run.bold = bold
    run.font.color.rgb = DocxRGBColor.from_string(color)


def _shade_cell(cell, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = properties.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        properties.append(shading)
    shading.set(qn("w:fill"), fill)


def _set_cell_margins(cell, *, top: int = 100, start: int = 120, bottom: int = 100, end: int = 120) -> None:
    properties = cell._tc.get_or_add_tcPr()
    margins = properties.first_child_found_in("w:tcMar")
    if margins is None:
        margins = OxmlElement("w:tcMar")
        properties.append(margins)
    for edge, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        tag = margins.find(qn(f"w:{edge}"))
        if tag is None:
            tag = OxmlElement(f"w:{edge}")
            margins.append(tag)
        tag.set(qn("w:w"), str(value))
        tag.set(qn("w:type"), "dxa")


def _set_table_geometry(table, widths_dxa: list[int], *, indent_dxa: int = 0) -> None:
    table.autofit = False
    properties = table._tbl.tblPr
    layout = properties.find(qn("w:tblLayout"))
    if layout is None:
        layout = OxmlElement("w:tblLayout")
        properties.append(layout)
    layout.set(qn("w:type"), "fixed")

    table_width = properties.find(qn("w:tblW"))
    if table_width is None:
        table_width = OxmlElement("w:tblW")
        properties.append(table_width)
    table_width.set(qn("w:w"), str(sum(widths_dxa)))
    table_width.set(qn("w:type"), "dxa")

    table_indent = properties.find(qn("w:tblInd"))
    if table_indent is None:
        table_indent = OxmlElement("w:tblInd")
        properties.append(table_indent)
    table_indent.set(qn("w:w"), str(indent_dxa))
    table_indent.set(qn("w:type"), "dxa")

    borders = properties.find(qn("w:tblBorders"))
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        properties.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        border = borders.find(qn(f"w:{edge}"))
        if border is None:
            border = OxmlElement(f"w:{edge}")
            borders.append(border)
        border.set(qn("w:val"), "single")
        border.set(qn("w:sz"), "5")
        border.set(qn("w:color"), LIGHT_LINE)

    grid = table._tbl.tblGrid
    for child in list(grid):
        grid.remove(child)
    for width in widths_dxa:
        column = OxmlElement("w:gridCol")
        column.set(qn("w:w"), str(width))
        grid.append(column)

    for row in table.rows:
        for index, cell in enumerate(row.cells):
            width = widths_dxa[index]
            cell.width = Inches(width / 1440)
            cell_properties = cell._tc.get_or_add_tcPr()
            cell_width = cell_properties.find(qn("w:tcW"))
            if cell_width is None:
                cell_width = OxmlElement("w:tcW")
                cell_properties.append(cell_width)
            cell_width.set(qn("w:w"), str(width))
            cell_width.set(qn("w:type"), "dxa")


def _add_response_line(document: Document) -> None:
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.space_before = Pt(3)
    paragraph.paragraph_format.space_after = Pt(8)
    paragraph.add_run(" ")
    properties = paragraph._p.get_or_add_pPr()
    borders = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")
    bottom.set(qn("w:space"), "2")
    bottom.set(qn("w:color"), LIGHT_LINE)
    borders.append(bottom)
    properties.append(borders)


def _configure_document(document: Document, *, worksheet: bool) -> None:
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)
    section.header_distance = Inches(0.35)
    section.footer_distance = Inches(0.35)

    styles = document.styles
    normal = styles["Normal"]
    normal.font.name = "Arial"
    normal._element.rPr.rFonts.set(qn("w:ascii"), "Arial")
    normal._element.rPr.rFonts.set(qn("w:hAnsi"), "Arial")
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(7 if worksheet else 6)
    normal.paragraph_format.line_spacing = 1.15

    for name, size, color, before, after in (
        ("Heading 1", 17, PURDUE_BLACK, 16, 7),
        ("Heading 2", 13, PURDUE_DARK_GOLD, 12, 5),
        ("Heading 3", 11.5, PURDUE_BLACK, 9, 4),
    ):
        style = styles[name]
        style.font.name = "Arial"
        style._element.rPr.rFonts.set(qn("w:ascii"), "Arial")
        style._element.rPr.rFonts.set(qn("w:hAnsi"), "Arial")
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = DocxRGBColor.from_string(color)
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)
        style.paragraph_format.keep_with_next = True

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.LEFT
    header.paragraph_format.space_after = Pt(0)
    _set_run_font(header.add_run("PURDUE UNIVERSITY"), size=9, bold=True, color=PURDUE_DARK_GOLD)
    _set_run_font(header.add_run("  |  COURSE DEVELOPMENT PARTNER"), size=9, color=MUTED)

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    footer.paragraph_format.space_before = Pt(0)
    _set_run_font(footer.add_run("Instructor-created teaching material"), size=8.5, color=MUTED)


def _add_doc_table(document: Document, table_spec: dict[str, Any]) -> None:
    headers = table_spec.get("headers") or []
    rows = table_spec.get("rows") or []
    if not headers or not rows:
        return
    table = document.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    usable_width_dxa = 9792
    base_width = usable_width_dxa // len(headers)
    widths_dxa = [base_width] * len(headers)
    widths_dxa[-1] += usable_width_dxa - sum(widths_dxa)
    for index, header in enumerate(headers):
        cell = table.rows[0].cells[index]
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        _shade_cell(cell, PURDUE_GOLD)
        _set_cell_margins(cell)
        cell.text = ""
        paragraph = cell.paragraphs[0]
        paragraph.paragraph_format.space_after = Pt(0)
        _set_run_font(paragraph.add_run(str(header)), size=10, bold=True, color=PURDUE_BLACK)
    for row_values in rows:
        cells = table.add_row().cells
        for index in range(len(headers)):
            value = str(row_values[index]) if index < len(row_values) else ""
            cell = cells[index]
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            _set_cell_margins(cell)
            cell.text = ""
            paragraph = cell.paragraphs[0]
            paragraph.paragraph_format.space_after = Pt(0)
            _set_run_font(paragraph.add_run(value), size=9.5)
    _set_table_geometry(table, widths_dxa)
    document.add_paragraph().paragraph_format.space_after = Pt(1)


def build_docx(spec: dict[str, Any], output_path: Path, project: dict[str, Any], *, worksheet: bool = False) -> None:
    document = Document()
    _configure_document(document, worksheet=worksheet)
    core = document.core_properties
    core.title = spec["title"]
    core.subject = "Course Development Partner teaching artifact"
    core.author = "Purdue University Course Development Partner"
    core.keywords = "education, course design, instructor authored"

    title = document.add_paragraph()
    title.paragraph_format.space_before = Pt(12)
    title.paragraph_format.space_after = Pt(4)
    _set_run_font(title.add_run(spec["title"]), size=27 if worksheet else 25, bold=True, color=PURDUE_BLACK)
    if spec.get("subtitle"):
        subtitle = document.add_paragraph()
        subtitle.paragraph_format.space_after = Pt(14)
        _set_run_font(subtitle.add_run(spec["subtitle"]), size=12, color=MUTED)
    context = document.add_paragraph()
    context.paragraph_format.space_after = Pt(14)
    label = "WORKSHEET" if worksheet else "TEACHING DOCUMENT"
    _set_run_font(context.add_run(label), size=9, bold=True, color=PURDUE_DARK_GOLD)
    course = project.get("course_name") or project.get("name") or "Course project"
    _set_run_font(context.add_run(f"  |  {course}"), size=9, color=MUTED)

    if worksheet:
        identity = document.add_paragraph()
        identity.paragraph_format.space_after = Pt(12)
        _set_run_font(identity.add_run("Name: ________________________________    Date: ________________"), size=10.5)

    for section in spec["sections"]:
        document.add_heading(section["heading"], level=1)
        if section.get("body"):
            paragraph = document.add_paragraph(section["body"])
            paragraph.paragraph_format.keep_together = False
        for bullet in section.get("bullets") or []:
            paragraph = document.add_paragraph(style="List Bullet")
            paragraph.paragraph_format.space_after = Pt(4)
            paragraph.add_run(bullet)
        for item in section.get("checklist") or []:
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.left_indent = Inches(0.15)
            paragraph.paragraph_format.first_line_indent = Inches(-0.15)
            _set_run_font(paragraph.add_run("☐  "), size=12, bold=True, color=PURDUE_DARK_GOLD)
            _set_run_font(paragraph.add_run(item), size=10.5)
        _add_doc_table(document, section.get("table") or {})
        for prompt in section.get("prompts") or []:
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.space_before = Pt(8)
            paragraph.paragraph_format.space_after = Pt(4)
            _set_run_font(paragraph.add_run(prompt), size=11, bold=True, color=PURDUE_BLACK)
            for _ in range(section.get("response_lines", 3)):
                _add_response_line(document)

    if spec.get("source_ids"):
        document.add_heading("Source trace", level=2)
        source_paragraph = document.add_paragraph()
        _set_run_font(source_paragraph.add_run("Project sources: "), size=9.5, bold=True)
        _set_run_font(source_paragraph.add_run(", ".join(spec["source_ids"])), size=9.5, color=MUTED)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    document.save(output_path)


def _set_pptx_text(paragraph, text: str, *, size: int, bold: bool = False, color: str = INK, font: str = "Arial") -> None:
    paragraph.text = text
    paragraph.font.name = font
    paragraph.font.size = PptxPt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = PptxRGBColor.from_string(color)


def _add_slide_notes(slide, source_ids: list[str]) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = "[Sources]\n" + ("\n".join(f"- Project source {source_id}" for source_id in source_ids) if source_ids else "- Instructor-provided project context; no external sources")


def build_pptx(spec: dict[str, Any], output_path: Path, project: dict[str, Any]) -> None:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)
    presentation.core_properties.title = spec["title"]
    presentation.core_properties.subject = "Course Development Partner slide deck"
    presentation.core_properties.author = "Purdue University Course Development Partner"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = PptxRGBColor.from_string(PURDUE_BLACK)
    title_box = title_slide.shapes.add_textbox(PptxInches(0.9), PptxInches(1.55), PptxInches(11.5), PptxInches(2.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_pptx_text(title_frame.paragraphs[0], spec["title"], size=52, bold=True, color="FFFFFF")
    subtitle = spec.get("subtitle") or project.get("course_name") or project.get("name") or "Course project"
    subtitle_box = title_slide.shapes.add_textbox(PptxInches(0.95), PptxInches(4.2), PptxInches(10.8), PptxInches(0.8))
    _set_pptx_text(subtitle_box.text_frame.paragraphs[0], subtitle, size=24, color=PURDUE_GOLD)
    brand_box = title_slide.shapes.add_textbox(PptxInches(0.95), PptxInches(6.55), PptxInches(5.5), PptxInches(0.35))
    _set_pptx_text(brand_box.text_frame.paragraphs[0], "PURDUE UNIVERSITY  |  COURSE DEVELOPMENT PARTNER", size=10, bold=True, color=PURDUE_GOLD)
    _add_slide_notes(title_slide, spec.get("source_ids") or [])

    for number, section in enumerate(spec["sections"], start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptxRGBColor.from_string("FFFFFF")
        heading_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.55), PptxInches(11.7), PptxInches(0.8))
        _set_pptx_text(heading_box.text_frame.paragraphs[0], section["heading"], size=36, bold=True, color=PURDUE_BLACK)

        body_box = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(1.65), PptxInches(11.25), PptxInches(4.95))
        frame = body_box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = PptxInches(0.06)
        frame.margin_right = PptxInches(0.06)
        frame.margin_top = PptxInches(0.06)
        frame.margin_bottom = PptxInches(0.06)
        content_items = []
        if section.get("body"):
            content_items.append((section["body"], False))
        content_items.extend((item, True) for item in (section.get("bullets") or []))
        content_items.extend((item, True) for item in (section.get("checklist") or []))
        content_items.extend((item, True) for item in (section.get("prompts") or []))
        if not content_items:
            content_items.append(("Instructor-guided discussion", False))
        for index, (text, bullet) in enumerate(content_items[:7]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            _set_pptx_text(paragraph, text, size=20 if bullet else 22, bold=False, color=INK)
            paragraph.level = 0
            paragraph.space_after = PptxPt(11)
            paragraph.line_spacing = 1.08
            if bullet:
                paragraph.text = f"•  {text}"

        footer = slide.shapes.add_textbox(PptxInches(0.85), PptxInches(7.02), PptxInches(11.7), PptxInches(0.25))
        footer_frame = footer.text_frame
        _set_pptx_text(footer_frame.paragraphs[0], f"PURDUE UNIVERSITY   •   {number}", size=9, bold=True, color=PURDUE_DARK_GOLD)
        footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _add_slide_notes(slide, spec.get("source_ids") or [])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def build_artifact_file(spec: dict[str, Any], output_path: Path, project: dict[str, Any]) -> None:
    kind = spec["kind"]
    if kind == "slides":
        build_pptx(spec, output_path, project)
    elif kind == "worksheet":
        build_docx(spec, output_path, project, worksheet=True)
    elif kind == "document":
        build_docx(spec, output_path, project, worksheet=False)
    else:
        raise ValueError(f"Unsupported artifact kind: {kind}")
